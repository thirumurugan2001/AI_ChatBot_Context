import psycopg2 
from dotenv import load_dotenv
from pptx import Presentation
import os
import psycopg2 
import pandas as pd
from azure.ai.inference import EmbeddingsClient
from azure.core.credentials import AzureKeyCredential
load_dotenv()

# Reads all the text from a PowerPoint (.pptx) file.
def read_pptx(file_path):
    prs = Presentation(file_path)
    slides_text = []

    for index, slide in enumerate(prs.slides):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text.append(shape.text.strip())
        slides_text.append(f"Slide {index + 1}:\n" + "\n".join(slide_text))
    
    return slides_text

# Function to get the embedding of the text
def get_embedding(text: str):
    try:
        token = os.getenv("AZURE_OPENAI_API_KEY")
        endpoint = os.getenv("AZURE_OPENAI_ENDPOINT")
        model_name = os.getenv("MODEL")
        client = EmbeddingsClient(
            endpoint=endpoint,
            credential=AzureKeyCredential(token)
        )
        response = client.embed(
            input=[text],
            model=model_name
        )
        embedding = response.data[0].embedding
        return embedding
    except Exception as e:
        print("Error in getting embedding: ", str(e))
        return None
    
# Function to connect to the database
def dbconnection():
    try:
        # Create a connection to the database
        connection = psycopg2.connect(
            dbname=os.getenv("DATABASE"),
            user=os.getenv("USER"),
            password= os.getenv("PASSWORD"),
            host= os.getenv("HOST"),
            port= os.getenv("PORT")
        )
        return connection
    except Exception as e:
        print(f"Error: Unable to connect to the database. {str(e)}")
        return None
    
# Function to select data into PostgreSQL
def selectData():
    try:
        conn = dbconnection()
        cursor = conn.cursor()
        cursor.execute("SELECT * FROM aboout_company")
        data = cursor.fetchall()
        return data
    except Exception as e:
        print("Error in selecting data: ", str(e))
        return None
    
# Function to select data into PostgreSQL
def dropTable():
    try:
        conn = dbconnection()
        cursor = conn.cursor()
        cursor.execute("DROP TABLE aboout_company")
        conn.commit()
        return "Table dropped successfully"
    except Exception as e:
        print("Error in selecting data: ", str(e))
        return None

# Function to create a table in PostgreSQL
def CreateTable():
    try:
        conn = dbconnection()
        if conn is not None:
            try :
                dropTable()
                print("Table is exist and droped successfully.")
            except Exception as e :
                print("Table does not exist, Continue to create table")
            cursor = conn.cursor()    
            cursor.execute("CREATE TABLE aboout_company (id SERIAL PRIMARY KEY,place TEXT NOT NULL,description TEXT,vectors VECTOR(3072));")
            conn.commit()
            return "Table created successfully"
        else:
            return "Failed to connect to the database"
    except Exception as e:
        print("Error in creating table: ", str(e))
        return None

# Function to read the excel file
def read_excel_file():
    try:
        file_path=r'Famous_Places_Tamil_Nadu.xlsx'
        df = pd.read_excel(file_path)
        print("Excel file read successfully")
        return df
    except Exception as e:
        print(f"Error reading the Excel file: {e}")
        return None

# Function to insert data into the database
def insert_into_database():
    try:
        conn = dbconnection()
        if conn is not None:
            df = read_excel_file()
            if df is not None:
                cursor = conn.cursor()
                for index, row in df.iterrows():
                    place = row['Place']
                    description = row['Description']
                    vector = get_embedding(description)
                    query = """INSERT INTO aboout_company (place, description, vectors) values (%s ,%s, %s)"""
                    cursor.execute(query, (place, description, vector))
                    conn.commit()
                    print(f"{index}th Row data inserted successfully")
                cursor.close()
                conn.close()
                return "Data inserted successfully"
            else:
                return "Error reading the Excel file"
        else :
            return "Connection failed"
    except Exception as e:
        print("Error in inserting data:", str(e))
        return None
